"""
Document processor
Handles PDF, DOCX, PPTX → text extraction → chunking → embedding → FAISS index
"""
import os
import re
import hashlib
import pickle
from pathlib import Path
from typing import List, Dict, Tuple

# ── Lazy imports (only loaded when needed) ───────────────────────────────────
def _import_pdf():
    try:
        import pdfplumber
        return pdfplumber
    except ImportError:
        import PyPDF2
        return PyPDF2

def _import_docx():
    from docx import Document
    return Document

def _import_pptx():
    from pptx import Presentation
    return Presentation

# ── Constants ────────────────────────────────────────────────────────────────
CHUNK_SIZE    = 800    # characters per chunk
CHUNK_OVERLAP = 150    # overlap between consecutive chunks
DATA_DIR      = Path("data")
UPLOAD_DIR    = DATA_DIR / "uploads"
INDEX_DIR     = DATA_DIR / "faiss_index"
DB_DIR        = DATA_DIR / "db"

for d in [UPLOAD_DIR, INDEX_DIR, DB_DIR]:
    d.mkdir(parents=True, exist_ok=True)

# ── Content type detection ───────────────────────────────────────────────────
CONTENT_TYPE_KEYWORDS = {
    "question_paper": [
        "question paper", "q.p.", "exam", "examination", "marks", "time:",
        "answer all", "section a", "section b", "part a", "part b",
        "unit i", "unit ii", "2 marks", "5 marks", "10 marks", "16 marks",
    ],
    "lab_manual": [
        "aim:", "objective:", "apparatus:", "procedure:", "result:",
        "experiment", "observation", "viva", "lab manual", "practical",
    ],
    "textbook": [
        "chapter", "definition", "theorem", "proof", "introduction",
        "summary", "exercises", "bibliography", "index",
    ],
    "lecture_notes": [
        "lecture", "notes", "topic:", "unit:", "module:", "slide",
    ],
}

def detect_content_type(text: str, filename: str) -> str:
    text_lower   = text[:2000].lower()
    fname_lower  = filename.lower()

    for ctype, keywords in CONTENT_TYPE_KEYWORDS.items():
        hits = sum(1 for kw in keywords if kw in text_lower or kw in fname_lower)
        if hits >= 2:
            return ctype

    # Fallback by filename hints
    for hint, ctype in [
        ("qp", "question_paper"), ("pyq", "question_paper"),
        ("lab", "lab_manual"),    ("exp", "lab_manual"),
        ("notes", "lecture_notes"), ("lec", "lecture_notes"),
        ("text", "textbook"),       ("book", "textbook"),
    ]:
        if hint in fname_lower:
            return ctype

    return "lecture_notes"

CONTENT_TYPE_LABELS = {
    "question_paper": "📝 Question Paper",
    "lab_manual":     "🔬 Lab Manual",
    "textbook":       "📖 Textbook",
    "lecture_notes":  "📋 Lecture Notes",
}

# ── Text extractors ──────────────────────────────────────────────────────────
def extract_pdf(filepath: str) -> str:
    text = ""
    try:
        import pdfplumber
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    text += t + "\n\n"
    except Exception:
        try:
            import PyPDF2
            with open(filepath, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    t = page.extract_text()
                    if t:
                        text += t + "\n\n"
        except Exception as e:
            text = f"[PDF extraction error: {e}]"
    return text.strip()


def extract_docx(filepath: str) -> str:
    try:
        from docx import Document
        doc  = Document(filepath)
        text = "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())

        # Also extract tables
        for table in doc.tables:
            rows = []
            for row in table.rows:
                rows.append(" | ".join(cell.text.strip() for cell in row.cells))
            text += "\n\n" + "\n".join(rows)

        return text.strip()
    except Exception as e:
        return f"[DOCX extraction error: {e}]"


def extract_pptx(filepath: str) -> str:
    try:
        from pptx import Presentation
        prs   = Presentation(filepath)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- Slide {i} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if len(slide_text) > 1:
                slides.append("\n".join(slide_text))
        return "\n\n".join(slides).strip()
    except Exception as e:
        return f"[PPTX extraction error: {e}]"


def extract_text(filepath: str) -> str:
    ext = Path(filepath).suffix.lower()
    if ext == ".pdf":
        return extract_pdf(filepath)
    elif ext in (".docx", ".doc"):
        return extract_docx(filepath)
    elif ext in (".pptx", ".ppt"):
        return extract_pptx(filepath)
    elif ext == ".txt":
        return Path(filepath).read_text(errors="ignore")
    else:
        return ""

# ── Chunker ──────────────────────────────────────────────────────────────────
def chunk_text(
    text: str,
    chunk_size: int = CHUNK_SIZE,
    overlap: int    = CHUNK_OVERLAP,
) -> List[str]:
    """Split text into overlapping chunks, respecting paragraph boundaries."""
    # Prefer paragraph-level splits
    paragraphs = re.split(r"\n{2,}", text)
    chunks, current = [], ""

    for para in paragraphs:
        para = para.strip()
        if not para:
            continue
        if len(current) + len(para) < chunk_size:
            current += ("\n\n" if current else "") + para
        else:
            if current:
                chunks.append(current)
            # If single paragraph is huge, hard-split it
            if len(para) > chunk_size:
                for i in range(0, len(para), chunk_size - overlap):
                    chunks.append(para[i : i + chunk_size])
            else:
                current = para

    if current:
        chunks.append(current)

    # Add overlap: each chunk carries the tail of the previous one
    overlapped = []
    for i, chunk in enumerate(chunks):
        if i > 0:
            tail = chunks[i - 1][-overlap:]
            chunk = tail + "\n\n" + chunk
        overlapped.append(chunk)

    return overlapped

# ── Document fingerprint ─────────────────────────────────────────────────────
def file_hash(filepath: str) -> str:
    h = hashlib.md5()
    with open(filepath, "rb") as f:
        for block in iter(lambda: f.read(65536), b""):
            h.update(block)
    return h.hexdigest()

# ── FAISS vector store helpers ────────────────────────────────────────────────
def build_vector_store(texts: List[str], metadatas: List[Dict], api_key: str):
    """Embed a list of text chunks and build a FAISS index."""
    from langchain_google_genai import GoogleGenerativeAIEmbeddings
    from langchain_community.vectorstores import FAISS

    os.environ["GOOGLE_API_KEY"] = api_key
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    store = FAISS.from_texts(texts, embeddings, metadatas=metadatas)
    return store


def save_vector_store(store, subject: str):
    path = INDEX_DIR / f"{subject.lower().replace(' ', '_')}.faiss"
    store.save_local(str(path))


def load_vector_store(subject: str, api_key: str):
    from langchain_google_genai import GoogleGenerativeAIEmbeddings
    from langchain_community.vectorstores import FAISS

    path = INDEX_DIR / f"{subject.lower().replace(' ', '_')}.faiss"
    if not path.exists():
        return None
    os.environ["GOOGLE_API_KEY"] = api_key
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    return FAISS.load_local(str(path), embeddings, allow_dangerous_deserialization=True)


def add_to_vector_store(existing_store, new_texts: List[str], new_meta: List[Dict], api_key: str):
    from langchain_google_genai import GoogleGenerativeAIEmbeddings
    from langchain_community.vectorstores import FAISS

    os.environ["GOOGLE_API_KEY"] = api_key
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    new_store  = FAISS.from_texts(new_texts, embeddings, metadatas=new_meta)
    if existing_store:
        existing_store.merge_from(new_store)
        return existing_store
    return new_store

# ── Main ingestion pipeline ──────────────────────────────────────────────────
def ingest_document(
    uploaded_file,
    subject: str,
    api_key: str,
    existing_store=None,
) -> Tuple[bool, str, dict]:
    """
    Full pipeline: save → extract → chunk → embed → store.
    Returns (success, message, metadata_dict).
    """
    if not api_key:
        return False, "❌ Please enter your Gemini API key first.", {}

    # Save file
    save_path = UPLOAD_DIR / uploaded_file.name
    save_path.write_bytes(uploaded_file.read())
    fhash = file_hash(str(save_path))

    # Extract text
    raw_text = extract_text(str(save_path))
    if not raw_text or raw_text.startswith("["):
        return False, f"❌ Could not extract text from {uploaded_file.name}.", {}

    # Detect content type
    content_type = detect_content_type(raw_text, uploaded_file.name)

    # Chunk
    chunks = chunk_text(raw_text)
    if not chunks:
        return False, "❌ No usable text found in document.", {}

    # Build metadata per chunk
    metadatas = [
        {
            "source":       uploaded_file.name,
            "subject":      subject,
            "content_type": content_type,
            "chunk_index":  i,
            "total_chunks": len(chunks),
            "file_hash":    fhash,
        }
        for i, _ in enumerate(chunks)
    ]

    # Embed + store
    try:
        store = add_to_vector_store(existing_store, chunks, metadatas, api_key)
    except Exception as e:
        return False, f"❌ Embedding error: {e}", {}

    doc_meta = {
        "filename":     uploaded_file.name,
        "subject":      subject,
        "content_type": content_type,
        "chunks":       len(chunks),
        "chars":        len(raw_text),
        "file_hash":    fhash,
    }

    label = CONTENT_TYPE_LABELS.get(content_type, content_type)
    return (
        True,
        f"✅ **{uploaded_file.name}** processed — {len(chunks)} chunks indexed as *{label}*.",
        doc_meta,
        store,
    )
